def extract_text(content: bytes, filename: str) -> str:
    ext = filename.lower().rsplit('.', 1)[-1] if '.' in filename else ''

    if ext == 'pdf':
        return _pdf(content)
    elif ext in ('pptx', 'ppt'):
        return _pptx(content)
    elif ext in ('docx', 'doc'):
        return _docx(content)
    elif ext in ('xlsx', 'xls'):
        return _xlsx(content)
    elif ext == 'csv':
        return content.decode('utf-8', errors='ignore')
    else:
        return content.decode('utf-8', errors='ignore')


def _pdf(content: bytes) -> str:
    import fitz
    doc = fitz.open(stream=content, filetype='pdf')
    return '\n'.join(page.get_text() for page in doc)


def _pptx(content: bytes) -> str:
    import io
    from pptx import Presentation
    prs = Presentation(io.BytesIO(content))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f'[Slide {i}]')
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                parts.append(shape.text.strip())
    return '\n'.join(parts)


def _docx(content: bytes) -> str:
    import io
    from docx import Document
    doc = Document(io.BytesIO(content))
    return '\n'.join(p.text for p in doc.paragraphs if p.text.strip())


def _xlsx(content: bytes) -> str:
    import io
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        parts.append(f'[Sheet: {sheet.title}]')
        for row in sheet.iter_rows(values_only=True):
            row_str = '\t'.join('' if c is None else str(c) for c in row)
            if row_str.strip():
                parts.append(row_str)
    return '\n'.join(parts)
